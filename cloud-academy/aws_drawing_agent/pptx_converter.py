"""
PPTX to React Flow Converter
=============================
Converts PowerPoint architecture diagrams to React Flow JSON format.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Optional, Tuple
import logging
import json
from pathlib import Path

logger = logging.getLogger(__name__)


class PPTXToReactFlowConverter:
    """
    Converts PowerPoint slides to React Flow diagram format.
    Extracts shapes, positions, connections, and text to create
    interactive diagrams for the platform.
    """
    
    # PowerPoint uses EMUs (English Metric Units)
    # 1 inch = 914400 EMUs
    # Standard slide: 10" x 7.5" = 9144000 x 6858000 EMUs
    EMU_PER_INCH = 914400
    
    def __init__(self, target_width: int = 1200, target_height: int = 800):
        """
        Initialize converter with target canvas dimensions.
        
        Args:
            target_width: Target canvas width in pixels
            target_height: Target canvas height in pixels
        """
        self.target_width = target_width
        self.target_height = target_height
    
    def convert_file(self, pptx_path: str) -> Dict:
        """
        Convert a PowerPoint file to React Flow format.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Dictionary with nodes and edges for React Flow
        """
        logger.info(f"Converting {pptx_path} to React Flow format")
        
        try:
            prs = Presentation(pptx_path)
            
            # Get slide dimensions for coordinate conversion
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Convert first slide (most architecture diagrams are single-slide)
            if len(prs.slides) == 0:
                logger.warning(f"No slides found in {pptx_path}")
                return {"nodes": [], "edges": []}
            
            slide = prs.slides[0]
            result = self._convert_slide(slide, slide_width, slide_height)
            
            # Add metadata
            result["metadata"] = {
                "source_file": Path(pptx_path).name,
                "slide_dimensions": {
                    "width": slide_width,
                    "height": slide_height,
                },
                "canvas_dimensions": {
                    "width": self.target_width,
                    "height": self.target_height,
                },
            }
            
            logger.info(f"Converted {len(result['nodes'])} nodes and {len(result['edges'])} edges")
            return result
            
        except Exception as e:
            logger.error(f"Error converting {pptx_path}: {e}")
            import traceback
            traceback.print_exc()
            return {"nodes": [], "edges": [], "error": str(e)}
    
    def _convert_slide(self, slide, slide_width: int, slide_height: int) -> Dict:
        """Convert a single slide to React Flow format."""
        nodes = []
        edges = []
        node_id_counter = 0
        
        # Process all shapes
        for shape in slide.shapes:
            result = self._process_shape(
                shape, 
                slide_width, 
                slide_height, 
                node_id_counter
            )
            
            if result:
                if result["type"] == "node":
                    nodes.append(result["data"])
                    node_id_counter += 1
                elif result["type"] == "edge":
                    edges.append(result["data"])
        
        return {
            "nodes": nodes,
            "edges": edges,
        }
    
    def _process_shape(
        self, 
        shape, 
        slide_width: int, 
        slide_height: int,
        node_id: int
    ) -> Optional[Dict]:
        """
        Process a single shape and convert to node or edge.
        
        Returns:
            Dict with type ("node" or "edge") and data, or None if shape should be skipped
        """
        # Skip certain shape types
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            return None
        
        # Handle LINE shapes as edges
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return self._convert_line_to_edge(shape, slide_width, slide_height)
        
        # Handle GROUP shapes recursively
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # For groups, we'll create a single node representing the group
            # In a more advanced version, we could process children separately
            pass
        
        # Convert shape to node
        node = self._convert_shape_to_node(shape, slide_width, slide_height, node_id)
        
        if node:
            return {"type": "node", "data": node}
        
        return None
    
    def _convert_shape_to_node(
        self, 
        shape, 
        slide_width: int, 
        slide_height: int,
        node_id: int
    ) -> Optional[Dict]:
        """Convert a shape to a React Flow node."""
        # Get position and size
        if not hasattr(shape, 'left'):
            return None
        
        # Convert coordinates from EMUs to pixels
        x, y = self._convert_coordinates(
            shape.left, 
            shape.top, 
            slide_width, 
            slide_height
        )
        
        width = int(shape.width / self.EMU_PER_INCH * 96)  # 96 DPI
        height = int(shape.height / self.EMU_PER_INCH * 96)
        
        # Extract text
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
        
        # Determine node type based on shape characteristics
        node_type = self._infer_node_type(shape, text)
        
        # Create node
        node = {
            "id": f"node-{node_id}",
            "type": node_type,
            "position": {"x": x, "y": y},
            "data": {
                "label": text,
                "width": width,
                "height": height,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type),
            },
        }
        
        # Add service identification if text matches AWS service
        service_id = self._identify_aws_service(text)
        if service_id:
            node["data"]["service_id"] = service_id
            node["type"] = "awsService"
        
        return node
    
    def _convert_line_to_edge(
        self, 
        shape, 
        slide_width: int, 
        slide_height: int
    ) -> Dict:
        """Convert a LINE shape to a React Flow edge."""
        # Lines in PowerPoint have start and end points
        # We'll need to find which nodes they connect
        
        # For now, create a basic edge structure
        # In production, we'd need to calculate which nodes this line connects
        edge = {
            "id": f"edge-{shape.name}",
            "source": "unknown",  # Would need to calculate
            "target": "unknown",  # Would need to calculate
            "type": "smoothstep",
            "data": {
                "shape_name": shape.name,
            },
        }
        
        return {"type": "edge", "data": edge}
    
    def _convert_coordinates(
        self, 
        emu_x: int, 
        emu_y: int, 
        slide_width: int, 
        slide_height: int
    ) -> Tuple[int, int]:
        """
        Convert PowerPoint EMU coordinates to React Flow pixel coordinates.
        
        Args:
            emu_x: X coordinate in EMUs
            emu_y: Y coordinate in EMUs
            slide_width: Slide width in EMUs
            slide_height: Slide height in EMUs
            
        Returns:
            Tuple of (x, y) in pixels
        """
        # Normalize to 0-1 range
        norm_x = emu_x / slide_width
        norm_y = emu_y / slide_height
        
        # Scale to target dimensions
        pixel_x = int(norm_x * self.target_width)
        pixel_y = int(norm_y * self.target_height)
        
        return pixel_x, pixel_y
    
    def _infer_node_type(self, shape, text: str) -> str:
        """
        Infer the React Flow node type based on shape characteristics.
        
        Returns:
            Node type string (e.g., "default", "awsService", "label", "group")
        """
        # Check if it's a large rectangle (likely a grouping/section)
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            width_inches = shape.width / self.EMU_PER_INCH
            height_inches = shape.height / self.EMU_PER_INCH
            
            if width_inches > 5 and height_inches > 3:
                return "group"
        
        # Check if it's a text box (label)
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return "label"
        
        # Check if it's a picture (likely an AWS icon)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "awsService"
        
        # Check if it's a group (likely service + icon)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "awsService"
        
        # Default node type
        return "default"
    
    def _identify_aws_service(self, text: str) -> Optional[str]:
        """
        Identify AWS service from text label.
        
        Args:
            text: Text from shape
            
        Returns:
            Service ID if identified, None otherwise
        """
        if not text:
            return None
        
        text_lower = text.lower().strip()
        
        # Common AWS service name mappings (must match frontend aws-services.ts IDs)
        service_mappings = {
            "ec2": "ec2",
            "elastic compute cloud": "ec2",
            "s3": "s3",
            "simple storage service": "s3",
            "lambda": "lambda",
            "aws lambda": "lambda",
            "rds": "rds",
            "relational database service": "rds",
            "dynamodb": "dynamodb",
            "vpc": "vpc",
            "virtual private cloud": "vpc",
            "iam": "iam",
            "cloudfront": "cloudfront",
            "route 53": "route53",
            "route53": "route53",
            "eks": "eks",
            "elastic kubernetes service": "eks",
            "ecs": "ecs",
            "elastic container service": "ecs",
            "api gateway": "api-gateway",
            "apigateway": "api-gateway",
            "amazon api gateway": "api-gateway",
            "cognito": "cognito",
            "sns": "sns",
            "sqs": "sqs",
            "cloudwatch": "cloudwatch",
            "cloudtrail": "cloudtrail",
            "kms": "kms",
            "secrets manager": "secrets-manager",
            "elasticache": "elasticache",
            "redshift": "redshift",
            "aurora": "aurora",
            "neptune": "neptune",
            "athena": "athena",
            "glue": "glue",
            "emr": "emr",
            "kinesis": "kinesis-streams",
            "kinesis data streams": "kinesis-streams",
            "kinesis firehose": "kinesis-firehose",
            "step functions": "step-functions",
            "stepfunctions": "step-functions",
            "aws step functions": "step-functions",
            "eventbridge": "eventbridge",
            "fargate": "fargate",
            "ecr": "ecr",
            "alb": "alb",
            "application load balancer": "alb",
            "nlb": "nlb",
            "network load balancer": "nlb",
            "elb": "alb",
            "elastic load balancer": "alb",
            "nat gateway": "nat-gateway",
            "internet gateway": "internet-gateway",
            "auto scaling": "auto-scaling",
            "auto scaling group": "auto-scaling",
            "waf": "waf",
            "shield": "shield",
            "guardduty": "guardduty",
            "opensearch": "opensearch",
            "elasticsearch": "opensearch",
            "msk": "msk",
            "kafka": "msk",
        }
        
        # Check for exact or partial matches
        for key, service_id in service_mappings.items():
            if key in text_lower:
                return service_id
        
        return None
    
    def convert_all_architectures(self, architectures_dir: str, output_dir: str):
        """
        Convert all PowerPoint files in a directory to React Flow format.
        
        Args:
            architectures_dir: Directory containing PPTX files
            output_dir: Directory to save JSON outputs
        """
        arch_path = Path(architectures_dir)
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        pptx_files = list(arch_path.glob("*.pptx"))
        logger.info(f"Converting {len(pptx_files)} architecture files")
        
        results = []
        for pptx_file in pptx_files:
            result = self.convert_file(str(pptx_file))
            
            # Save to JSON
            output_file = output_path / f"{pptx_file.stem}.json"
            with open(output_file, 'w') as f:
                json.dump(result, f, indent=2)
            
            results.append({
                "file": pptx_file.name,
                "nodes": len(result.get("nodes", [])),
                "edges": len(result.get("edges", [])),
                "output": str(output_file),
            })
        
        logger.info(f"Converted {len(results)} files")
        return results
