#!/usr/bin/env python3
"""
Analyze PowerPoint file structure to understand how to extract drawing elements.
This will help us understand how AWS architecture diagrams are structured.
"""

import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from pathlib import Path


def analyze_shape(shape, depth=0):
    """Analyze a single shape and return its properties."""
    indent = "  " * depth
    shape_info = {
        "name": shape.name,
        "type": str(shape.shape_type),
        "has_text": shape.has_text_frame,
    }
    
    # Get position and size
    if hasattr(shape, 'left'):
        shape_info["position"] = {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
        }
    
    # Get text content
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            shape_info["text"] = text
    
    # Check if it's a picture/image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        shape_info["is_image"] = True
        if hasattr(shape, 'image'):
            shape_info["image_type"] = shape.image.content_type
    
    # Check if it's a group (contains other shapes)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        shape_info["is_group"] = True
        shape_info["children"] = []
        for child_shape in shape.shapes:
            child_info = analyze_shape(child_shape, depth + 1)
            shape_info["children"].append(child_info)
    
    # Check for connectors/lines
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        shape_info["is_connector"] = True
    
    # Auto shape (rectangles, circles, etc.)
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        shape_info["is_auto_shape"] = True
        if hasattr(shape, 'auto_shape_type'):
            shape_info["auto_shape_type"] = str(shape.auto_shape_type)
    
    return shape_info


def analyze_slide(slide, slide_num):
    """Analyze a single slide and extract all shapes."""
    print(f"\n{'='*60}")
    print(f"SLIDE {slide_num}")
    print(f"{'='*60}")
    
    slide_info = {
        "slide_number": slide_num,
        "shapes": []
    }
    
    print(f"Total shapes on slide: {len(slide.shapes)}")
    
    for idx, shape in enumerate(slide.shapes):
        shape_info = analyze_shape(shape)
        slide_info["shapes"].append(shape_info)
        
        print(f"\nShape {idx + 1}:")
        print(f"  Name: {shape_info['name']}")
        print(f"  Type: {shape_info['type']}")
        
        if "position" in shape_info:
            pos = shape_info["position"]
            print(f"  Position: ({pos['left']}, {pos['top']})")
            print(f"  Size: {pos['width']} x {pos['height']}")
        
        if "text" in shape_info:
            print(f"  Text: {shape_info['text'][:100]}...")
        
        if shape_info.get("is_image"):
            print(f"  Image Type: {shape_info.get('image_type', 'unknown')}")
        
        if shape_info.get("is_group"):
            print(f"  Group with {len(shape_info['children'])} children")
        
        if shape_info.get("is_connector"):
            print(f"  Connector/Line")
    
    return slide_info


def analyze_pptx(file_path):
    """Analyze a PowerPoint file and extract structure."""
    print(f"\nAnalyzing: {file_path}")
    print(f"{'='*80}\n")
    
    try:
        prs = Presentation(file_path)
        
        print(f"Presentation Info:")
        print(f"  Total Slides: {len(prs.slides)}")
        print(f"  Slide Width: {prs.slide_width}")
        print(f"  Slide Height: {prs.slide_height}")
        
        analysis = {
            "file": str(file_path),
            "total_slides": len(prs.slides),
            "slide_dimensions": {
                "width": prs.slide_width,
                "height": prs.slide_height,
            },
            "slides": []
        }
        
        # Analyze first 3 slides (or all if less than 3)
        max_slides = min(3, len(prs.slides))
        for idx in range(max_slides):
            slide = prs.slides[idx]
            slide_info = analyze_slide(slide, idx + 1)
            analysis["slides"].append(slide_info)
        
        # Save analysis to JSON
        output_file = Path(file_path).stem + "_analysis.json"
        with open(output_file, 'w') as f:
            json.dump(analysis, f, indent=2, default=str)
        
        print(f"\n{'='*80}")
        print(f"Analysis saved to: {output_file}")
        print(f"{'='*80}\n")
        
        return analysis
        
    except Exception as e:
        print(f"Error analyzing {file_path}: {e}")
        import traceback
        traceback.print_exc()
        return None


def main():
    """Main function."""
    if len(sys.argv) < 2:
        print("Usage: python analyze_pptx_structure.py <path_to_pptx_file>")
        print("\nExample:")
        print("  python analyze_pptx_structure.py /path/to/diagram.pptx")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    if not Path(file_path).exists():
        print(f"Error: File not found: {file_path}")
        sys.exit(1)
    
    analysis = analyze_pptx(file_path)
    
    if analysis:
        print("\nKey Findings:")
        print(f"  - Total slides: {analysis['total_slides']}")
        print(f"  - Slide dimensions: {analysis['slide_dimensions']['width']} x {analysis['slide_dimensions']['height']}")
        
        total_shapes = sum(len(slide['shapes']) for slide in analysis['slides'])
        print(f"  - Total shapes analyzed: {total_shapes}")
        
        # Count shape types
        shape_types = {}
        for slide in analysis['slides']:
            for shape in slide['shapes']:
                shape_type = shape['type']
                shape_types[shape_type] = shape_types.get(shape_type, 0) + 1
        
        print(f"\n  Shape type distribution:")
        for shape_type, count in sorted(shape_types.items(), key=lambda x: x[1], reverse=True):
            print(f"    - {shape_type}: {count}")


if __name__ == "__main__":
    main()
